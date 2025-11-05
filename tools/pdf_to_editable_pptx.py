#!/usr/bin/env python3
"""
Convert PDF to editable PPTX with exact layout preservation.

This tool uses a hybrid approach:
1. Renders PDF pages as high-resolution background images (300 DPI)
2. Overlays transparent text boxes with exact positioning and styling
3. Results in PPTX that looks identical to original PDF + fully editable text

The generated PPTX maintains:
- Exact visual layout (via background image)
- Fully editable text (via transparent text boxes)
- Complete styling (fonts, colors, sizes)
- Reasonable file size (2-3MB with background images)

Usage:
    python tools/pdf_to_editable_pptx.py input.pdf
    python tools/pdf_to_editable_pptx.py input.pdf output.pptx

Output:
    By default, saves to examples/ folder as {filename}_editable.pptx

Requirements:
    pip install pymupdf python-pptx pillow
"""

import sys
import argparse
from pathlib import Path
from typing import List, Dict, Tuple
import logging
import io

try:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from PIL import Image
except ImportError as e:
    print(f"❌ Missing required package: {e}")
    print("Install with: pip install pymupdf python-pptx pillow")
    sys.exit(1)

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

# Constants
DPI = 300  # Resolution for background image rendering
MARGIN_INCHES = 0.01  # Minimal margin for text boxes
PDF_TO_INCHES = 1 / 72  # PDF points to inches conversion


class PDFToEditablePPTXHybrid:
    """Convert PDF to editable PPTX with exact layout preservation using hybrid method."""

    def __init__(self):
        """Initialize converter."""
        self.doc = None
        self.prs = None

    def pdf_to_emu(self, pdf_value: float) -> int:
        """Convert PDF points to EMU (English Metric Units)."""
        return int(pdf_value * 914400 / 72)

    def extract_page_data(self, pdf_path: str) -> List[Dict]:
        """
        Extract PDF pages with text content and rendering data.

        Args:
            pdf_path: Path to PDF file

        Returns:
            List of page dictionaries with background image and text data
        """
        pages_data = []

        try:
            self.doc = fitz.open(pdf_path)
            total = len(self.doc)

            for page_num, page in enumerate(self.doc, 1):
                print(f"   Extracting: {page_num}/{total} pages", end='\r')

                # Render page as high-resolution image for background
                pix = page.get_pixmap(matrix=fitz.Matrix(DPI / 72, DPI / 72))
                img_bytes = pix.tobytes("png")

                # Extract text with positioning
                text_data = page.get_text("dict")
                text_spans = []

                if "blocks" in text_data:
                    for block in text_data["blocks"]:
                        if block.get("type") == 0:  # Text block
                            for line in block.get("lines", []):
                                for span in line.get("spans", []):
                                    span_info = {
                                        'text': span.get("text", ""),
                                        'font': span.get("font", "Helvetica"),
                                        'size': span.get("size", 12),
                                        'color': span.get("color", 0),
                                        'flags': span.get("flags", 0),
                                        'bbox': span.get("bbox", (0, 0, 0, 0)),
                                    }
                                    text_spans.append(span_info)

                page_data = {
                    'page_num': page_num,
                    'width': page.rect.width,
                    'height': page.rect.height,
                    'image_bytes': img_bytes,
                    'image_width': pix.width,
                    'image_height': pix.height,
                    'text_spans': text_spans,
                }

                pages_data.append(page_data)

            print()  # New line after progress
            logger.info(f"✅ Extracted {total} pages with background images")

        except Exception as e:
            logger.error(f"❌ Error extracting PDF: {e}")
            raise
        finally:
            if self.doc:
                self.doc.close()

        return pages_data

    def create_hybrid_pptx(
        self,
        pages_data: List[Dict],
        output_path: str
    ) -> None:
        """
        Create PPTX with background images and overlay text boxes.

        Args:
            pages_data: List of page data with images and text
            output_path: Path to output PPTX file
        """
        logger.info("📐 Creating hybrid PPTX (background + text overlay)...")

        self.prs = Presentation()

        # Set slide dimensions to match PDF aspect ratio exactly
        first_page = pages_data[0]
        pdf_width_points = first_page['width']
        pdf_height_points = first_page['height']

        # Convert PDF points to inches (1 point = 1/72 inch)
        pdf_width_inches = pdf_width_points / 72
        pdf_height_inches = pdf_height_points / 72

        self.prs.slide_width = Inches(pdf_width_inches)
        self.prs.slide_height = Inches(pdf_height_inches)

        for page_data in pages_data:
            # Create blank slide
            blank_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_layout)

            # Add background image to fill slide
            img_stream = io.BytesIO(page_data['image_bytes'])
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(
                img_stream,
                left, top,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )

            # Group text spans by lines
            lines = self._group_spans_by_line(page_data['text_spans'])

            # Add transparent text boxes over background
            for line_spans in lines:
                if not line_spans:
                    continue

                try:
                    self._add_text_line(slide, line_spans, page_data)
                except Exception as e:
                    logger.warning(f"⚠️  Failed to add text line: {e}")

            # Progress indicator
            progress = int((page_data['page_num'] / len(pages_data)) * 100)
            print(
                f"   Creating: {progress}% ({page_data['page_num']}/{len(pages_data)})",
                end='\r'
            )

        # Save PPTX
        try:
            self.prs.save(output_path)
            print()  # New line after progress
            logger.info(f"✅ Successfully created: {output_path}")
            logger.info(f"   Total slides: {len(pages_data)}")
        except Exception as e:
            print()  # New line after progress
            logger.error(f"❌ Error saving PPTX: {e}")
            raise

    def _group_spans_by_line(
        self,
        spans: List[Dict],
        threshold: float = 3.0
    ) -> List[List[Dict]]:
        """Group spans by Y position (lines) with improved logic."""
        if not spans:
            return []

        # Sort by Y coordinate, then by X coordinate
        sorted_spans = sorted(spans, key=lambda s: (s['bbox'][1], s['bbox'][0]))
        lines = []

        if not sorted_spans:
            return lines

        current_line = [sorted_spans[0]]
        current_y = sorted_spans[0]['bbox'][1]

        for span in sorted_spans[1:]:
            span_y = span['bbox'][1]
            span_height = span['bbox'][3] - span['bbox'][1]

            # Use adaptive threshold based on font size
            font_size = span.get('size', 12)
            adaptive_threshold = min(threshold, font_size * 0.3)

            # Check if span belongs to current line
            if abs(current_y - span_y) <= adaptive_threshold:
                current_line.append(span)
            else:
                # Sort current line by X coordinate before saving
                current_line.sort(key=lambda s: s['bbox'][0])
                lines.append(current_line)

                # Start new line
                current_line = [span]
                current_y = span_y

        # Add the last line
        if current_line:
            current_line.sort(key=lambda s: s['bbox'][0])
            lines.append(current_line)

        return lines

    def _add_text_line(
        self,
        slide,
        line_spans: List[Dict],
        page_data: Dict
    ) -> None:
        """
        Add a text line as transparent text box overlaying background.

        Args:
            slide: PPTX slide
            line_spans: List of spans forming a line
            page_data: Page data with dimensions
        """
        # Get bounding box for entire line
        first_span = line_spans[0]
        last_span = line_spans[-1]

        x0, y0, x1, y1 = first_span['bbox']
        _, _, last_x1, _ = last_span['bbox']

        # Map PDF coordinates directly to slide coordinates (1:1 ratio)
        slide_height_points = page_data['height']
        left = Inches(x0 / 72)
        top = Inches((slide_height_points - y0) / 72)
        width = Inches((last_x1 - x0) / 72)
        height = Inches((y1 - y0) / 72)

        # Ensure minimum dimensions
        width = max(width, Inches(0.5))
        height = max(height, Inches(0.1))

        # Create text box with transparent background
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = False
        text_frame.margin_bottom = Inches(MARGIN_INCHES)
        text_frame.margin_left = Inches(MARGIN_INCHES)
        text_frame.margin_right = Inches(MARGIN_INCHES)
        text_frame.margin_top = Inches(MARGIN_INCHES)

        # Make text box transparent and remove border
        textbox.fill.background()
        textbox.line.color.rgb = RGBColor(0, 0, 0)
        textbox.line.width = 0

        # Clear default paragraph text
        text_frame.paragraphs[0].text = ""

        # Add text with styling using proper TextRun management
        paragraph = text_frame.paragraphs[0]

        for span_idx, span in enumerate(line_spans):
            text = span['text'].strip()
            if not text:
                continue

            # Clean text
            clean_text = self._clean_text_for_xml(text)
            if not clean_text:
                continue

            # Create or get run for this span
            if span_idx == 0:
                # Use the first run of the paragraph
                if len(paragraph.runs) > 0:
                    run = paragraph.runs[0]
                    run.text = clean_text
                else:
                    run = paragraph.add_run()
                    run.text = clean_text
            else:
                # Create new run for additional spans
                run = paragraph.add_run()
                run.text = clean_text

            # Apply styling to this specific run
            self._apply_span_style(run, span)

    def _apply_span_style(self, run, span: Dict) -> None:
        """Apply font styling from span to text run."""
        try:
            font = run.font

            # Font name
            font_name = self._normalize_font_name(span.get('font', 'Helvetica'))
            font.name = font_name

            # Font size
            font.size = Pt(span.get('size', 12))

            # Color
            color_int = span.get('color', 0)
            if color_int is not None and color_int != 0:
                r, g, b = self._color_int_to_rgb(color_int)
                font.color.rgb = RGBColor(r, g, b)

            # Bold and italic
            flags = span.get('flags', 0)
            font.bold = bool(flags & (1 << 3))
            font.italic = bool(flags & (1 << 0))

        except Exception as e:
            logger.warning(f"⚠️  Failed to apply style: {e}")

    def _color_int_to_rgb(self, color_int: int) -> Tuple[int, int, int]:
        """Convert PDF color integer to RGB tuple."""
        if color_int is None or color_int == 0:
            return (0, 0, 0)

        r = (color_int >> 16) & 0xFF
        g = (color_int >> 8) & 0xFF
        b = color_int & 0xFF

        return (r, g, b)

    def _normalize_font_name(self, font_name: str) -> str:
        """Normalize PDF font name to PowerPoint-compatible name."""
        name_map = {
            'Helv': 'Helvetica',
            'Times': 'Times New Roman',
            'Courier': 'Courier New',
        }

        for pdf_name, pptx_name in name_map.items():
            if pdf_name.lower() in font_name.lower():
                return pptx_name

        return font_name or 'Helvetica'

    def _clean_text_for_xml(self, text: str) -> str:
        """Clean text to ensure XML compatibility."""
        if not text:
            return ""

        # XML 1.0不允许的控制字符
        illegal_chars = [
            # C0控制字符（除了tab, newline, carriage return）
            c for c in range(0x00, 0x20) if c not in (0x09, 0x0A, 0x0D)
        ] + [
            # C1控制字符
            c for c in range(0x80, 0x9F)
        ] + [
            # 其他XML无效字符
            0xFFFE, 0xFFFF
        ]

        # 清理文本
        cleaned = []
        for char in text:
            char_code = ord(char)

            # 跳过非法字符
            if char_code in illegal_chars:
                continue
            elif char in '&<>':
                # XML特殊字符需要转义
                if char == '&':
                    cleaned.append('&amp;')
                elif char == '<':
                    cleaned.append('&lt;')
                elif char == '>':
                    cleaned.append('&gt;')
            else:
                cleaned.append(char)

        return ''.join(cleaned)

    def convert(self, pdf_path: str, output_path: str) -> None:
        """
        Convert PDF to editable PPTX with layout preservation.

        Args:
            pdf_path: Path to input PDF file
            output_path: Path to output PPTX file
        """
        pdf_path = Path(pdf_path)
        output_path = Path(output_path)

        if not pdf_path.exists():
            logger.error(f"❌ PDF file not found: {pdf_path}")
            sys.exit(1)

        logger.info(f"📄 Converting PDF to editable PPTX: {pdf_path.name}")
        logger.info(f"   (Hybrid: background image + editable text overlay)")

        # Extract pages with images and text
        pages_data = self.extract_page_data(str(pdf_path))

        # Create hybrid PPTX
        self.create_hybrid_pptx(pages_data, str(output_path))


def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description='Convert PDF to editable PPTX with exact layout preservation',
        epilog='Example: python pdf_to_editable_pptx.py input.pdf'
    )
    parser.add_argument('input', help='Input PDF file path')
    parser.add_argument(
        'output',
        nargs='?',
        help='Output PPTX file path (default: examples/{filename}_editable.pptx)'
    )

    args = parser.parse_args()

    # Default output path to examples folder
    if args.output is None:
        input_path = Path(args.input)
        output_path = Path('examples') / f"{input_path.stem}_editable.pptx"
    else:
        output_path = Path(args.output)

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Convert PDF to PPTX
    converter = PDFToEditablePPTXHybrid()
    try:
        converter.convert(args.input, str(output_path))
    except Exception as e:
        logger.error(f"❌ Conversion failed: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
