#!/usr/bin/env python3
"""
Convert PDF to PPTX presentation, preserving layout by converting each slide to an image.

Usage:
    python tools/pdf_to_pptx.py input.pdf output.pptx [--dpi DPI] [--add-notes]

Requirements:
    pip install pdf2image python-pptx pillow

System dependencies:
    # Linux/WSL
    sudo apt install poppler-utils

    # macOS
    brew install poppler
"""

import sys
import argparse
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io
import os


def pdf_to_pptx(pdf_path, output_path, dpi=300, add_notes=False):
    """
    Convert PDF to PPTX by embedding each page as an image.

    Args:
        pdf_path (str): Path to input PDF file
        output_path (str): Path to output PPTX file
        dpi (int): Resolution for PDF conversion (default: 300)
        add_notes (bool): Add page numbers as notes (default: False)
    """
    pdf_path = Path(pdf_path)
    output_path = Path(output_path)

    if not pdf_path.exists():
        print(f"Error: PDF file not found: {pdf_path}")
        sys.exit(1)

    print(f"📄 Converting PDF to PPTX: {pdf_path.name}")
    print(f"   DPI: {dpi}, Add notes: {add_notes}")

    # Convert PDF pages to images
    try:
        print(f"🔄 Extracting PDF pages...")
        images = convert_from_path(str(pdf_path), dpi=dpi)
        print(f"✅ Extracted {len(images)} pages")
    except Exception as e:
        print(f"❌ Error converting PDF: {e}")
        print("   Make sure poppler-utils is installed:")
        print("   Linux/WSL: sudo apt install poppler-utils")
        print("   macOS: brew install poppler")
        sys.exit(1)

    # Create PPTX presentation
    print(f"📐 Creating PPTX presentation...")
    prs = Presentation()

    # Standard slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Process each image
    for idx, image in enumerate(images, 1):
        # Convert PIL Image to bytes
        img_bytes = io.BytesIO()
        image.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Calculate scaling to fit slide while preserving aspect ratio
        img_width, img_height = image.size
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Calculate scale factor
        width_ratio = slide_width / Inches(img_width / dpi)
        height_ratio = slide_height / Inches(img_height / dpi)
        scale = min(width_ratio, height_ratio)

        pic_width = Inches(img_width / dpi * scale)
        pic_height = Inches(img_height / dpi * scale)

        # Center image on slide
        left = (slide_width - pic_width) / 2
        top = (slide_height - pic_height) / 2

        # Add image to slide
        slide.shapes.add_picture(
            img_bytes,
            left, top,
            width=pic_width,
            height=pic_height
        )

        # Add notes if requested
        if add_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = f"Page {idx}/{len(images)}"

        # Progress indicator
        progress = int((idx / len(images)) * 100)
        print(f"   Processing: {progress}% ({idx}/{len(images)})", end='\r')

    # Save PPTX
    try:
        prs.save(output_path)
        print(f"\n✅ Successfully created: {output_path}")
        print(f"   Total slides: {len(images)}")
    except Exception as e:
        print(f"\n❌ Error saving PPTX: {e}")
        sys.exit(1)


def main():
    parser = argparse.ArgumentParser(
        description='Convert PDF to PPTX presentation',
        epilog='Example: python pdf_to_pptx.py slides.pdf output.pptx --dpi 300'
    )
    parser.add_argument('input', help='Input PDF file path')
    parser.add_argument('output', nargs='?', help='Output PPTX file path (default: input with .pptx extension)')
    parser.add_argument('--dpi', type=int, default=300, help='DPI for PDF conversion (default: 300)')
    parser.add_argument('--add-notes', action='store_true', help='Add page numbers as slide notes')

    args = parser.parse_args()

    # Default output path
    if args.output is None:
        input_path = Path(args.input)
        args.output = input_path.with_suffix('.pptx')

    pdf_to_pptx(args.input, args.output, dpi=args.dpi, add_notes=args.add_notes)


if __name__ == '__main__':
    main()
