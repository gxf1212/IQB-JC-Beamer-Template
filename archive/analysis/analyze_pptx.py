#!/usr/bin/env python3
"""分析 PPTX 文件的布局结构和内容"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
from pathlib import Path


def analyze_slide(slide, slide_num):
    """分析单个幻灯片的布局"""
    print(f"\n{'='*60}")
    print(f"幻灯片 {slide_num}")
    print(f"{'='*60}")

    # 布局名称
    print(f"布局类型: {slide.slide_layout.name}")

    # 统计元素
    shapes_info = {
        'text': 0,
        'picture': 0,
        'table': 0,
        'chart': 0,
        'shape': 0,
        'group': 0,
        'other': 0
    }

    images = []
    texts = []

    for shape in slide.shapes:
        # 统计形状类型
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_info['picture'] += 1
            # 获取图片位置和大小信息
            images.append({
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'name': shape.name
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shapes_info['table'] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            shapes_info['chart'] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shapes_info['group'] += 1
        elif hasattr(shape, 'text') and shape.text.strip():
            shapes_info['text'] += 1
            texts.append({
                'text': shape.text.strip()[:50],  # 前50个字符
                'left': shape.left,
                'top': shape.top,
                'font_size': getattr(shape.text_frame.paragraphs[0].runs[0].font, 'size', None) if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs else None
            })
        else:
            shapes_info['other'] += 1

    # 打印统计
    print(f"\n元素统计:")
    for shape_type, count in shapes_info.items():
        if count > 0:
            print(f"  - {shape_type}: {count}")

    # 打印图片布局
    if images:
        print(f"\n图片布局 (共 {len(images)} 张):")
        for i, img in enumerate(images, 1):
            print(f"  图片 {i}: 位置({img['left']//914400:.1f}cm, {img['top']//914400:.1f}cm), "
                  f"大小({img['width']//914400:.1f}cm × {img['height']//914400:.1f}cm)")

    # 打印文本框
    if texts:
        print(f"\n文本内容 (共 {len(texts)} 个):")
        for i, txt in enumerate(texts[:5], 1):  # 只显示前5个
            size_info = f", 字号:{txt['font_size']//12700}pt" if txt['font_size'] else ""
            print(f"  文本 {i}: \"{txt['text']}...\"{size_info}")

    return shapes_info, images, texts


def analyze_presentation(pptx_path):
    """分析整个演示文稿"""
    print(f"\n{'#'*60}")
    print(f"分析文件: {pptx_path}")
    print(f"{'#'*60}")

    try:
        prs = Presentation(pptx_path)

        # 幻灯片尺寸
        print(f"\n幻灯片尺寸: {prs.slide_width//914400}cm × {prs.slide_height//914400}cm")
        print(f"总幻灯片数: {len(prs.slides)}")

        # 分析每一页
        all_stats = []
        for i, slide in enumerate(prs.slides, 1):
            stats, imgs, txts = analyze_slide(slide, i)
            all_stats.append({
                'slide_num': i,
                'stats': stats,
                'images': imgs,
                'texts': txts
            })

        # 总结
        print(f"\n{'='*60}")
        print("总体布局特点总结:")
        print(f"{'='*60}")

        # 统计常见布局模式
        total_images = sum(len(s['images']) for s in all_stats)
        total_texts = sum(len(s['texts']) for s in all_stats)
        avg_images = total_images / len(all_stats)
        avg_texts = total_texts / len(all_stats)

        print(f"平均每页图片数: {avg_images:.1f}")
        print(f"平均每页文本框数: {avg_texts:.1f}")

        # 识别布局模式
        print(f"\n常见布局模式:")
        for s in all_stats:
            img_count = len(s['images'])
            txt_count = len(s['texts'])
            if img_count == 0 and txt_count > 0:
                pattern = "纯文本页"
            elif img_count == 1 and txt_count > 0:
                pattern = "单图+文字"
            elif img_count == 2:
                pattern = "双图并列"
            elif img_count >= 3:
                pattern = f"{img_count}图并列"
            elif img_count > 0 and txt_count == 0:
                pattern = "纯图片"
            else:
                pattern = "其他"

            print(f"  幻灯片 {s['slide_num']}: {pattern}")

        return all_stats

    except Exception as e:
        print(f"错误: {e}")
        import traceback
        traceback.print_exc()
        return None


if __name__ == "__main__":
    # 分析多个文件
    pptx_files = [
        "/mnt/e/GitHub-repo/literature-reading/JC/2025.6/25.6-xufan-JC.pptx",
        "/mnt/e/GitHub-repo/literature-reading/JC/2024.12.17/24.12.17-xufan-JC.pptx",
        "/mnt/e/GitHub-repo/literature-reading/JC/2024.7/24.7.10-xufan-JC.pptx",
        "/mnt/e/GitHub-repo/literature-reading/JC/2023.11/23.11.8-xufan-JC.pptx",
    ]

    for pptx_file in pptx_files:
        if Path(pptx_file).exists():
            analyze_presentation(pptx_file)
        else:
            print(f"\n文件不存在: {pptx_file}")
